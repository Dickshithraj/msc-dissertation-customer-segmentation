"""Build docs/presentation.pptx — 25-minute supervisor presentation."""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent.parent  # repo root (scripts/ is one level down)
FIG = ROOT / "outputs" / "figures"
OUT = ROOT / "docs" / "presentation.pptx"

NAVY = RGBColor(0x1F, 0x4E, 0x79)
BLUE = RGBColor(0x2E, 0x74, 0xB5)
TEAL = RGBColor(0x11, 0x7A, 0x82)
GREY = RGBColor(0x40, 0x40, 0x40)
LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
GOLD = RGBColor(0xBF, 0x8F, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

SW, SH = 13.333, 7.5


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, color, line=False):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if line:
        sh.line.color.rgb = color
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def textbox(slide, x, y, w, h, lines, size=18, color=GREY, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.12,
            font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, item in enumerate(lines):
        # item: str  OR  (text, overrides-dict)
        if isinstance(item, tuple):
            text, ov = item
        else:
            text, ov = item, {}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ov.get("align", align)
        p.line_spacing = line_spacing
        p.space_after = Pt(ov.get("space_after", 6))
        run = p.add_run()
        run.text = text
        f = run.font
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        f.italic = ov.get("italic", False)
        f.color.rgb = ov.get("color", color)
        f.name = ov.get("font", font)
        if ov.get("bullet"):
            pPr = p._pPr if p._pPr is not None else p.get_or_add_pPr()
    return tb


def header(slide, title, tag=None):
    rect(slide, 0, 0, SW, 0.16, TEAL)
    textbox(slide, 0.55, 0.32, 11.0, 0.9, title, size=30, color=NAVY, bold=True)
    if tag:
        textbox(slide, 10.4, 0.42, 2.6, 0.5, tag, size=14, color=GOLD, bold=True,
                align=PP_ALIGN.RIGHT)
    ln = rect(slide, 0.55, 1.12, 12.2, 0.025, BLUE)


def bullets(slide, items, x=0.55, y=1.45, w=6.6, h=5.6, size=17):
    lines = []
    for it in items:
        if isinstance(it, tuple):
            text, ov = it
        else:
            text, ov = it, {}
        prefix = "" if ov.pop("noprefix", False) else "•  "
        indent = ov.pop("indent", 0)
        lines.append(((" " * 5 * indent) + prefix + text, ov))
    textbox(slide, x, y, w, h, lines, size=size)


def picture(slide, path, x, y, w=None, h=None):
    kw = {}
    if w:
        kw["width"] = Inches(w)
    if h:
        kw["height"] = Inches(h)
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), **kw)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def stat_row(slide, stats, y=5.9, x0=0.55, box_w=2.9, gap=0.25, h=1.1):
    x = x0
    for big, small in stats:
        rect(slide, x, y, box_w, h, LIGHT)
        textbox(slide, x + 0.1, y + 0.08, box_w - 0.2, 0.55, big, size=21,
                color=TEAL, bold=True, align=PP_ALIGN.CENTER)
        textbox(slide, x + 0.1, y + 0.6, box_w - 0.2, 0.5, small, size=11.5,
                color=GREY, align=PP_ALIGN.CENTER)
        x += box_w + gap


# ── 1. Title ─────────────────────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, 4.85, SW, 0.06, TEAL)
textbox(s, 1.0, 2.0, 11.3, 1.9,
        "Optimizing Targeted Marketing:\nA Framework for Customer Segmentation and "
        "Automated Cluster-Based Notification",
        size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
textbox(s, 1.0, 5.1, 11.3, 1.6, [
    ("Dickshith Raj Nagaraj", {"size": 22, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}),
    ("MSc Advanced Computer Science — University of Leeds — Session 2025/2026",
     {"size": 16, "color": RGBColor(0xBD, 0xD7, 0xEE), "align": PP_ALIGN.CENTER}),
    ("Supervisor: Dr Fahad Panolan",
     {"size": 16, "color": RGBColor(0xBD, 0xD7, 0xEE), "align": PP_ALIGN.CENTER}),
])
notes(s, "Open with the one-line story: raw transactions go in at the top; a costed, "
         "per-customer marketing action comes out at the bottom — automatically, "
         "reproducibly, and with the money quantified BEFORE any campaign is sent.")

# ── 2. Problem ───────────────────────────────────────────────────────────────
s = add_slide()
header(s, "The Problem: Blanket Marketing", "2 min")
bullets(s, [
    ("Retailers hold rich transactional data …", {"size": 20}),
    ("… yet most outbound marketing is still “blanket”: every customer gets the same message.",
     {"size": 20}),
    ("Doubly wasteful:", {"size": 20, "bold": True}),
    ("spends contact budget on customers who will never respond", {"indent": 1, "size": 18}),
    ("fatigues high-value customers with irrelevant contact — eroding the very "
     "lifetime value it should protect", {"indent": 1, "size": 18}),
], w=12.2)
rect(s, 0.55, 5.3, 12.2, 1.5, LIGHT)
textbox(s, 0.85, 5.5, 11.6, 1.2,
        "Research question: how do you convert a raw retail transaction log — automatically, "
        "reproducibly, and with quantified financial justification — into a per-customer, "
        "prioritised marketing-notification plan?",
        size=18, color=NAVY, bold=True)
notes(s, "Keep this to 90 seconds. The framed box is the thesis of the whole talk.")

# ── 3. Aim & objectives ─────────────────────────────────────────────────────
s = add_slide()
header(s, "Aim & Six Objectives", "2 min")
bullets(s, [
    ("O1 — Background research: systematic survey justifying every method choice", {}),
    ("O2 — Data foundation: auditable cleaning + RFM+ feature engineering", {}),
    ("O3 — Validated segmentation: 6 algorithms benchmarked, pre-declared selection rule", {}),
    ("O4 — Predictive layer: probabilistic CLV + leakage-guarded churn models", {}),
    ("O5 — Decision layer: rule-based notification engine + dashboard", {}),
    ("O6 — Commercial quantification: Monte Carlo ROI vs untargeted baseline", {}),
], w=12.2, size=19)
stat_row(s, [("UCI Online Retail II", "public dataset, 2009–2011"),
             ("~160 s", "full pipeline, commodity CPU"),
             ("Stages 1 → 10", "one auditable pipeline"),
             ("Exploratory software", "MSc project type")])
notes(s, "One breath per objective. Emphasise: each objective is measurable and each is "
         "discharged in an identified chapter of the report.")

# ── 4. Architecture ─────────────────────────────────────────────────────────
s = add_slide()
header(s, "System Architecture — Five Layers, Stages 1–10", "2 min")
picture(s, FIG / "system_architecture.png", 8.15, 1.25, h=6.1)
bullets(s, [
    ("Layer 1 — Data foundation: Stage 1 load & clean, 2 features, 2b preprocessing", {}),
    ("Layer 2 — Segmentation: Stage 3 clustering ×6, 3b validation, 4 profiling", {}),
    ("Layer 3 — Customer intelligence: Stage 5 CLV, 6 churn, 7 migration", {}),
    ("Layer 4 — Decision engine: Stage 8 notifications, 9 Monte Carlo ROI", {}),
    ("Layer 5 — Delivery: Stage 10 dashboard + artefact store", {}),
    ("Data flows strictly one way; every stage persists its artefacts → "
     "independently re-runnable, trivially auditable", {"bold": True, "color": NAVY}),
], w=7.3, size=16)
notes(s, "Point down the figure as you name the layers. Close with the one-way data flow "
         "point — it is what makes the system auditable.")

# ── 5. Data & cleaning ──────────────────────────────────────────────────────
s = add_slide()
header(s, "Stages 1–2b — Data, Cleaning, Features", "3 min")
bullets(s, [
    ("Four sequential, logged cleaning rules (every exclusion measured):", {"bold": True}),
    ("missing Customer ID  −243,007 rows (22.77%)", {"indent": 1, "size": 16}),
    ("cancellations (“C” invoices)  −18,744 (2.27%)", {"indent": 1, "size": 16}),
    ("invalid quantity/price  −89 (0.01%)", {"indent": 1, "size": 16}),
    ("exact duplicates  −11,940 (1.48%)", {"indent": 1, "size": 16}),
    ("7 behavioural features per customer: RFM + Tenure, AvgOrderValue, "
     "AvgInterPurchaseDays, DistinctProducts", {}),
    ("Skew-gated preprocessing: log1p only where |skew| > 0.5, then z-scaling — "
     "correct geometry for distance-based clustering", {}),
], w=7.3, size=17)
picture(s, FIG / "scaling_effect.png", 8.0, 1.7, w=5.1)
stat_row(s, [("1.07M → 794K", "raw → audited rows (−25.65%)"),
             ("5,878", "unique customers"),
             ("7", "behavioural features"),
             ("Audit table", "persisted for every rule")])
notes(s, "Stress the audit table: nothing is silently dropped — this is the evidence an "
         "examiner or auditor would ask for.")

# ── 6. Clustering ───────────────────────────────────────────────────────────
s = add_slide()
header(s, "Stage 3 — Six Clustering Algorithms, Five Families", "4 min")
picture(s, FIG / "cluster_pca_projection.png", 8.35, 1.3, h=5.9)
bullets(s, [
    ("Literature: no single algorithm dominates → choose empirically per dataset", {"bold": True}),
    ("K-Means · GMM · DBSCAN · HDBSCAN · Agglomerative (Ward) · Spectral", {}),
    ("Automatic hyper-parameter selection:", {}),
    ("K-Means: elbow (k=4) vs silhouette (k=2) — heuristics disagree; reported, not hidden",
     {"indent": 1, "size": 15}),
    ("GMM: BIC picks 10 components — a warning sign", {"indent": 1, "size": 15}),
    ("DBSCAN: ε from k-distance knee · HDBSCAN: min cluster size 50 (~1% of base)",
     {"indent": 1, "size": 15}),
    ("All six projected on identical PCA axes (76.6% variance) for like-for-like comparison", {}),
], w=7.5, size=16)
notes(s, "This and the next slide are your strongest methodology minutes. The point is not "
         "any single algorithm — it is that the choice is made by evidence.")

# ── 7. Validation & selection ───────────────────────────────────────────────
s = add_slide()
header(s, "Stage 3b — Validation, Stability, Pre-Declared Selection", "")
picture(s, FIG / "stability_ari.png", 7.9, 1.35, w=5.25)
bullets(s, [
    ("Three internal indices: silhouette, Davies–Bouldin, Calinski–Harabasz", {}),
    ("Stability: 50 bootstrap resamples, Adjusted Rand Index vs full-data partition", {}),
    ("Selection rule (pre-declared, applied mechanically):", {"bold": True, "color": NAVY}),
    ("1. disqualify noise fraction > 30%", {"indent": 1, "size": 16}),
    ("2. require mean bootstrap ARI ≥ 0.70", {"indent": 1, "size": 16}),
    ("3. among survivors: highest silhouette", {"indent": 1, "size": 16}),
    ("Winner: HDBSCAN — silhouette 0.416, best overall rank 1.75", {"bold": True, "color": TEAL}),
    ("Honest negatives: DBSCAN collapsed (ARI ≈ 0, eliminated mechanically); "
     "GMM over-fragmented (silhouette 0.081); K-Means default rejected by evidence", {"size": 15}),
], w=7.1, size=16)
notes(s, "Say the rule verbatim — being pre-declared is what makes the selection "
         "reproducible and criticisable. Rigour changed the answer: K-Means was rejected.")

# ── 8. Segments ─────────────────────────────────────────────────────────────
s = add_slide()
header(s, "Stage 4 — Three Operational Segments", "")
picture(s, FIG / "segment_profiles.png", 7.55, 1.45, w=5.6)
bullets(s, [
    ("General Customers — 3,218 (54.8%): the mainstream; moderate recency, "
     "frequency ≈7, spend ≈£2,422", {}),
    ("Lost Customers — 1,546 (26.3%): exactly one purchase, tenure 0, "
     "recency 363 days — a sharply defined lapsed group", {}),
    ("Noise / Uncategorised — 1,114 (19.0%): a finding, not a failure — "
     "the most valuable, most heterogeneous customers (mean spend £8,467)", {"bold": True}),
    ("Deterministic naming rules → same profile always gets the same marketing name", {}),
    ("Noise customers are excluded from generic campaigns and flagged for individual "
     "treatment — marketing-meaningful, ethically safer", {"color": NAVY}),
], w=6.8, size=17)
notes(s, "The counter-intuitive highlight: the customers HDBSCAN refuses to classify are "
         "the highest-value accounts. Forcing them into a segment would mis-target them.")

# ── 9. CLV ──────────────────────────────────────────────────────────────────
s = add_slide()
header(s, "Stage 5 — Customer Lifetime Value (BG/NBD + Gamma–Gamma)", "4 min")
picture(s, FIG / "clv_distribution.png", 7.3, 1.6, w=5.8)
bullets(s, [
    ("Probabilistic models built for non-contractual retail — customers lapse "
     "silently, never “cancel”", {}),
    ("Independence assumption (frequency vs value) checked before fitting: ≈ 0 ✓", {}),
    ("Outputs per customer: P(alive), expected purchases at 90/180/365 days, "
     "12-month discounted CLV", {}),
    ("Classic long tail: top decile holds a disproportionate share of value", {}),
    ("Role in the system: relative value ranking that feeds the CLV tiers of the "
     "decision engine", {"color": NAVY}),
], w=6.6, size=17)
stat_row(s, [("£8.33M", "12-month portfolio CLV"),
             ("£371 / £1,418", "median / mean CLV"),
             ("0.909", "mean P(alive)"),
             ("90/180/365 d", "purchase forecasts")])
notes(s, "Honest caveat if asked: no temporal hold-out (data ends 2011), so CLV is used as "
         "a ranking, not a validated forecast — listed in threats to validity.")

# ── 10. Churn ───────────────────────────────────────────────────────────────
s = add_slide()
header(s, "Stage 6 — Churn: Six Classifiers Under a Leakage Guard", "")
picture(s, FIG / "churn_roc_curves.png", 7.9, 1.35, w=5.2)
bullets(s, [
    ("Label: Recency > 90th percentile (535 days) → 10% churners", {}),
    ("Leakage guard: label is derived from Recency → Recency EXCLUDED from features "
     "(verified by unit test)", {"bold": True, "color": NAVY}),
    ("Best: Random Forest — ROC-AUC 0.849, recall 0.752 …", {}),
    ("… but McNemar exact test: RF vs XGBoost p = 0.878 → statistically tied; "
     "no over-claimed “winner”", {}),
    ("Deliberate contrast: GB & KNN left unweighted → ROC-AUC ≈0.85 but recall "
     "collapses to 0.03–0.07 — ROC-AUC alone misleads under imbalance", {"bold": True}),
    ("Error asymmetry: a missed churner costs £371+ CLV; a wasted retention email "
     "costs pence → recall-heavy operating points are economically rational", {"size": 15}),
], w=7.1, size=16)
notes(s, "“An AUC near 1.0 would have been a red flag for leakage; 0.849 with the "
         "guard is a credible ceiling.” Drivers: Tenure 0.37, cadence 0.22.")

# ── 11. Notification engine ─────────────────────────────────────────────────
s = add_slide()
header(s, "Stage 8 — Rule-Based Notification Engine", "3 min")
bullets(s, [
    ("Where analysis becomes action — the gap most published studies leave open", {"bold": True}),
    ("Inputs per customer: segment name + CLV tier (80th/40th pct) + churn band (0.50/0.25)", {}),
    ("Step 1 — segment playbook: action, channel, offer, base priority (1–5)", {}),
    ("Step 2 — value & risk modulation: high-risk + valuable → “Priority retention "
     "intervention” on premium channels; low-value high-risk → cheap automated reactivation", {}),
    ("Step 3 — contact timing: 0.6 × the customer’s own inter-purchase interval "
     "— no fixed drumbeat, over-contact suppressed by construction", {}),
    ("Why rules, not a learned policy? Auditability (NFR5): every recommendation "
     "traceable, inspectable, overridable by a human", {"color": NAVY, "bold": True}),
], w=12.2, size=17)
stat_row(s, [("5,878", "campaigns planned (all customers)"),
             ("64  (1.1%)", "priority retention interventions"),
             ("3,023", "engagement promotions"),
             ("recommend(id)", "on-demand API → dashboard")])
notes(s, "Read one plan row aloud from the dashboard in the demo: segment, CLV, risk, "
         "action, priority, contact day. Expensive attention is rare by construction.")

# ── 12. ROI ─────────────────────────────────────────────────────────────────
s = add_slide()
header(s, "Stage 9 — Monte Carlo ROI vs Blanket Baseline", "2 min")
picture(s, FIG / "roi_distribution.png", 7.35, 1.5, w=5.75)
bullets(s, [
    ("10,000 simulated worlds: Beta response priors → Binomial conversions → "
     "noisy revenue − real channel costs", {}),
    ("Targeted plan: mean ROI 38.4×, 95% CI [21.8×, 60.9×], "
     "P(ROI>0) = 100%", {"bold": True}),
    ("Static blanket baseline (same machinery): 30.4×", {}),
    ("Headline: +168% profit uplift (+£15,025)", {"bold": True, "color": TEAL, "size": 19}),
    ("Why the uplift is the defensible number: both plans share the same priors, margin "
     "and cost model → assumption errors cancel in the difference", {"color": NAVY}),
], w=6.6, size=16)
stat_row(s, [("£623.50", "targeted contact cost"),
             ("£293.90", "blanket contact cost"),
             ("38.4× vs 30.4×", "targeted vs blanket ROI"),
             ("+168%", "profit uplift")], y=6.15, h=1.0)
notes(s, "Pre-empt the objection yourself: response rates are documented planning "
         "assumptions, not fitted — which is exactly why the uplift, not the absolute "
         "ROI, is the headline. A/B calibration is future work item #1.")

# ── 13. Dashboard ───────────────────────────────────────────────────────────
s = add_slide()
header(s, "Stage 10 — Interactive Dashboard  (live demo)", "3 min")
picture(s, FIG / "dashboard_overview.png", 6.7, 1.45, w=6.3)
bullets(s, [
    ("Eight Streamlit pages reading only persisted artefacts", {}),
    ("Overview · Segments · Lifetime Value · Churn Risk · Migration · "
     "Notifications · ROI · Customer Lookup", {"size": 15}),
    ("Demo route: Overview → Segments → ROI → finish on Customer Lookup", {"bold": True}),
    ("Lookup answers, for any Customer ID: “what should we send this person, "
     "and why?” — with the why traceable to explicit rules", {"color": NAVY}),
    ("Turns the batch pipeline into a tool a non-technical marketer can operate", {}),
], w=6.0, size=16)
notes(s, "Have Streamlit already running before the meeting. If the demo fails, this "
         "screenshot and the PDF figures are the fallback.")

# ── 14. Rigour ──────────────────────────────────────────────────────────────
s = add_slide()
header(s, "Engineering Rigour & Reproducibility", "2 min")
bullets(s, [
    ("Fixed seed (42) + pinned snapshot date → byte-identical outputs on every run", {}),
    ("20 pytest unit tests on synthetic data — cleaning rules, leakage guard, "
     "priority cap, ROI arithmetic all asserted", {}),
    ("Single configuration module — no magic numbers; the config file IS the "
     "experiment registry", {}),
    ("Artefact-mediated coupling: stages communicate only through files → any stage "
     "re-runnable and inspectable in isolation", {}),
    ("≈160 s end-to-end on commodity CPU — no GPU, no cluster", {}),
    ("Ethics by design: pseudonymised public data, GDPR principles, noise customers get "
     "the LEAST aggressive action, cadence-based timing discourages over-contact", {}),
], w=12.2, size=18)
notes(s, "One breath for the whole slide if time is short — but do not skip the tests "
         "and reproducibility lines; they anchor the professional-practice marks.")

# ── 15. Limitations & future work ───────────────────────────────────────────
s = add_slide()
header(s, "Honest Limitations & Future Work", "")
bullets(s, [
    ("Limitations (volunteered, not hidden):", {"bold": True, "color": NAVY}),
    ("single dataset (one UK retailer, 2009–2011) — no generalisation claim", {"indent": 1, "size": 16}),
    ("coarse structure (2 clusters + noise) — what the evidence supports; forcing more degrades quality", {"indent": 1, "size": 16}),
    ("CLV: no temporal hold-out → used as ranking, not validated forecast", {"indent": 1, "size": 16}),
    ("ROI priors are planning assumptions → hence the uplift framing", {"indent": 1, "size": 16}),
    ("Future work (in order of value):", {"bold": True, "color": NAVY}),
    ("A/B-calibrated uplift modelling — target by incremental response", {"indent": 1, "size": 16}),
    ("temporal validation of CLV & churn on a longer window", {"indent": 1, "size": 16}),
    ("real-time deployment: recommend() behind a service API on streaming features", {"indent": 1, "size": 16}),
    ("learned embeddings + GPU implementations at larger scale — the validation "
     "framework transfers unchanged", {"indent": 1, "size": 16}),
], w=12.2, size=17)
notes(s, "Volunteering limitations earns marks. Each limitation pairs with a future-work "
         "item that would resolve it.")

# ── 16. Conclusions ─────────────────────────────────────────────────────────
s = add_slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, 1.7, SW, 0.05, TEAL)
textbox(s, 0.9, 0.55, 11.5, 1.0, "Conclusions", size=34, color=WHITE, bold=True)
textbox(s, 0.9, 2.1, 11.5, 3.6, [
    ("1.  Rigour changes answers — the evidence rejected the K-Means default, "
     "blocked an unjustified “best model” claim, and exposed two silently "
     "failing churn models.", {"size": 20, "color": WHITE, "space_after": 14}),
    ("2.  The closed loop is the contribution — validated clusters, value forecasts "
     "and risk scores only become worth computing when a traceable decision layer turns "
     "them into costed actions.", {"size": 20, "color": WHITE, "space_after": 14}),
    ("3.  The whole loop fits in one auditable, three-minute, commodity-hardware "
     "pipeline — raw transactions in, a costed action per customer out.",
     {"size": 20, "color": WHITE}),
])
textbox(s, 0.9, 6.1, 11.5, 0.9, "Thank you — questions welcome",
        size=22, color=RGBColor(0xBD, 0xD7, 0xEE), bold=True, align=PP_ALIGN.CENTER)
notes(s, "Close on the loop: most published studies stop at the clusters; this project "
         "goes all the way to the campaign plan and the money.")

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print("wrote", OUT)
